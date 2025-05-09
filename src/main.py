import os
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFacePipeline, HuggingFaceEmbeddings
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_text_splitters import RecursiveCharacterTextSplitter
from transformers import AutoModelForCausalLM, AutoTokenizer, pipeline

# Document Processing
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        return " ".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        return " ".join(page.extract_text() for page in reader.pages if page.extract_text())
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

# def extract_text_from_pptx(file_path):
#     try:
#         prs = Presentation(file_path)
#         text = []
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text") and shape.text.strip():
#                     text.append(shape.text)
#         return " ".join(text)
#     except Exception as e:
#         print(f"Error reading {file_path}: {e}")
#         return ""

def extract_all_text(directory):
    all_text = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        if filename.endswith(".docx"):
            all_text.append(extract_text_from_docx(file_path))
        elif filename.endswith(".pdf"):
            all_text.append(extract_text_from_pdf(file_path))
        # elif filename.endswith(".pptx"):
        #     all_text.append(extract_text_from_pptx(file_path))
    return " ".join(all_text)

# Text Chunking
def chunk_text(all_text):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    chunks = text_splitter.split_text(text=all_text)
    print(f"Created {len(chunks)} chunks")
    return chunks

# Vector Store with Semantic Embeddings
def create_vector_store(chunks):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    vector_store = FAISS.from_texts(chunks, embeddings)
    return vector_store

# Format Retrieved Documents
def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

# RAG Chain
def generate_answer(question, vector_store, history=None):
    # Initialize local LLM (CPU-optimized)
    model_id = "distilgpt2"  # Lightweight, CPU-friendly model
    tokenizer = AutoTokenizer.from_pretrained(model_id)
    model = AutoModelForCausalLM.from_pretrained(model_id)
    pipe = pipeline(
        "text-generation",
        model=model,
        tokenizer=tokenizer,
        max_new_tokens=100,
        device=-1  # Force CPU
    )
    llm = HuggingFacePipeline(pipeline=pipe)

    # Define retriever for semantic search
    retriever = vector_store.as_retriever(search_kwargs={"k": 4})

    # Define prompt
    prompt = ChatPromptTemplate.from_messages([
        ("system", "Answer based on the following context: {context}\nIf you don't know, say so."),
        MessagesPlaceholder(variable_name="history"),
        ("human", "{question}")
    ])

    # Create RAG chain
    rag_chain = (
        {"context": retriever | format_docs, "question": RunnablePassthrough(), "history": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    # Invoke with history
    history = history or []
    return rag_chain.invoke({"question": question, "history": history})

# Main Execution
def main():
    # Extract text
    directory = "data"  # Adjust to your document folder
    all_text = extract_all_text(directory)
    if not all_text:
        print("No text extracted. Check document directory.")
        return

    # Chunk text
    chunks = chunk_text(all_text)
    if not chunks:
        print("No chunks created.")
        return

    # Create vector store with semantic embeddings
    vector_store = create_vector_store(chunks)

    # Example query
    question = "What is the main topic of the documents?"
    response = generate_answer(question, vector_store)
    print(f"Query: {question}\nResponse: {response}")

if __name__ == "__main__":
    main()