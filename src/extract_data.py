import os
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation


# Document Processing
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        return " ".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        return " ".join(page.extract_text() for page in reader.pages if page.extract_text())
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return " ".join(text)
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_all_text(directory):
    all_text = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        if filename.endswith(".docx"):
            all_text.append(extract_text_from_docx(file_path))
        elif filename.endswith(".pdf"):
            all_text.append(extract_text_from_pdf(file_path))
        elif filename.endswith(".pptx"):
            all_text.append(extract_text_from_pptx(file_path))
    return " ".join(all_text)



